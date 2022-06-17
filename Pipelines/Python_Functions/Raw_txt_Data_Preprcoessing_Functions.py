# Extract all audios from video
# Convert all the text transcript from audios
# return the size of the video and extracted audio files
# Extract the topics and subtopics from meta data PPT
# Organize the data into the labeld data format in tabular format

import speech_recognition as sr 
import moviepy.editor as mp
import pandas as pd
import numpy as np
import glob
import os


list_of_videos = glob.glob('../../Data/videos_univ_chicago/*.mp4')
list_of_ppts =  glob.glob('../../Data/text_univ_chichago/*.pptx')

list_of_ppts = [x for x in list_of_ppts if "~$" not in x]


# get video files size on the video in Bytes
def file_size(vid_file):
    import os
    byts = os.path.getsize(vid_file)
    # then we return the bytes into bytes, metabyte, gb
    return byts, np.round(byts/1000000, 3), np.round(byts/1000000000, 3)



# get video files durations 
def video_duration(vid_file):
    from moviepy.editor import VideoFileClip
    clip = VideoFileClip(vid_file)
    return clip.duration,  np.round(clip.duration/60, 3)


# Function to read all the video clips and convert to audio files 
def convert_to_audio(vd,  output_pth):
    clip = mp.VideoFileClip(vd) 
    ad_nm = vd.split("/")[-1].split(".")[0]
    clip.audio.write_audiofile(output_pth + "/" + ad_nm + ".wav")
    print("Audio {} has been successfully extracted in {}".format(ad_nm, output_pth))


# convert all video clips into Audio files 
# for vid in list_of_videos:
#     convert_to_audio(vd = vid, output_pth = 'Audios')



# Reading Large Audio Files
# If you want to perform speech recognition of a long audio file, then the below function handles that quite well:
# https://www.thepythoncode.com/article/using-speech-recognition-to-convert-speech-to-text-python

# read all the audo files from the path
# Convert audio into text files 

r = sr.Recognizer()
audio = sr.AudioFile("converted.wav")

# importing libraries 
import speech_recognition as sr 
import os 
from pydub import AudioSegment 
from pydub.silence import split_on_silence 

# create a speech recognition object
r = sr.Recognizer()

# a function that splits the audio file into chunks
# and applies speech recognition
def get_large_audio_transcription(path):
    """
    Splitting the large audio file into chunks
    and apply speech recognition on each of these chunks
    """
    # open the audio file using pydub
    sound = AudioSegment.from_wav(path)  
    # split audio sound where silence is 700 miliseconds or more and get chunks
    chunks = split_on_silence(sound,
        # experiment with this value for your target audio file
        min_silence_len = 500,
        # adjust this per requirement
        silence_thresh = sound.dBFS-14,
        # keep the silence for 1 second, adjustable as well
        keep_silence=500,
    )
    
    folder_name = "audio-chunks"
    # create a directory to store the audio chunks
    if not os.path.isdir(folder_name):
        os.mkdir(folder_name)
    whole_text = ""
    # process each chunk 
    for i, audio_chunk in enumerate(chunks, start=1):
        # export audio chunk and save it in
        # the `folder_name` directory.
        chunk_filename = os.path.join(folder_name, f"chunk{i}.wav")
        audio_chunk.export(chunk_filename, format="wav")
        # recognize the chunk
        with sr.AudioFile(chunk_filename) as source:
            audio_listened = r.record(source)
            # try converting it to text
            try:
                text = r.recognize_google(audio_listened)
            except sr.UnknownValueError as e:
                print("Error:", str(e))
            else:
                text = f"{text.capitalize()}. "
                print(chunk_filename, ":", text)
                whole_text += text
    fl_nm = path.split("/")[-1].split(".")[0]
    res_df = pd.DataFrame({"video_file" : [str(fl_nm)] ,  "raw_txt": [str(whole_text)] } )
    # return the text for all chunks detected and a dataframe that contain the text chunk
    return whole_text, res_df

# convert the convertsed txt into .txt files 
def txt_file(raw_txt,  text_output):
    with open(text_output, mode ='w') as file: 
        file.write("Recognized Speech:") 
        file.write("\n") 
        file.write(raw_txt) 
        print("text ready!")


# audio_rawtxt = pd.DataFrame()
# for audio in list_of_audios:
#     print(audio + "\n")
#     extracted_txt = get_large_audio_transcription(path = audio)
#     txt_file(raw_txt = extracted_txt[0], text_output = "Audio_Text/" + path.split("/")[-1].replace(".wav", ".txt") )
#     temp = extracted_txt[1]
#     audio_rawtxt = pd.concat([audio_rawtxt, temp], axis = 0)
# audio_rawtxt.to_csv("Targeted_Output/Phase1/0_raw_audio_rawtxt.csv", index = False)




# Convert meta data ppt into text files
# https://stackoverflow.com/questions/13559133/how-to-open-ppt-file-using-python
# https://python-pptx.readthedocs.io/en/latest/api/shapes.html#shape-objects-autoshapes

# function to extract all text from the ppt slides
def read_ppt_text(ppt_path):
    from pptx import Presentation
    prs = Presentation(ppt_path)

    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
    #         print(shape)
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs


# ppts_names = [i.split("/")[-1].split(".")[0] for i in list_of_ppts]
# meta_data_lst = [read_ppt_text(ppt_path = i) for i in list_of_ppts]

# metadata_dict = dict(zip(ppts_names, meta_data_lst))
# raw_metadata = pd.DataFrame(metadata_dict.items(), columns = ['meta_file', "raw_text"])
# raw_metadata.head()
# raw_metadata.to_csv("Targeted_Output/Phase1/0_raw_metadata.csv", index = False)




# Annotate audio files with timestamp
# https://towardsdatascience.com/speech-recognition-with-timestamps-934ede4234b2

# the way that work:

# split each audio in each 5-10 second wav.files while keep track of the second durations in each files incrementally
# read each audio file and extract the text
# and mark each text extracted from each splitted audio file with the incremntally tracked time stamps

# inspared by https://stackoverflow.com/questions/37999150/how-to-split-a-wav-file-into-multiple-wav-files

# inspared by https://stackoverflow.com/questions/37999150/how-to-split-a-wav-file-into-multiple-wav-files

from pydub import AudioSegment
import math
import re

# function to split audios into different samller clips 
def split_audios_clips(folder, filename, clip_length_sec = 5 ):

    audio = AudioSegment.from_wav(folder + filename)

    # devide the auido by each 5 secs ==> 46 clips
    total_secs = math.ceil(audio.duration_seconds / clip_length_sec)

    sec_per_split = 1 # increment by 1
    start_time = 0
    start_time_lst = []
    end_time_lst = []

    for i in range(0, total_secs, sec_per_split):

        split_fn = str(i) + '_' + filename

        from_sec = i
        to_sec = i+sec_per_split
        t1 = from_sec * clip_length_sec * 1000
        t2 = to_sec * clip_length_sec * 1000

        split_audio = audio[t1:t2]
        split_audio.export(folder + '/splitted_audio/' + split_fn, format="wav")

        start_time_lst.append(start_time)
        start_time = start_time + clip_length_sec
        end_time_lst.append(start_time)
        
    return start_time_lst, end_time_lst

# function to split audio file into smaller files and then extract sentences from each small clips and uniformaly generate time stamp
def generate_timestamped_sentences(audio_path, temp_outpt_wac = 'Audios/splitted_audio/', clip_length_sec = 10 ):
    
    folder = audio_path.split("/")[0]  + "/"
    filename = audio_path.split("/")[1]  
    filetype = audio_path.split("/")[1].split(".")[0]
    
    # divide the audio into clips by 10 seconda each 
    time_stamps = split_audios_clips(clip_length_sec = clip_length_sec, folder = folder, filename = filename)
    
    # save start and end sec time
    time_stamps_df = pd.DataFrame()
    time_stamps_df["start sec"] = time_stamps[0]
    time_stamps_df["end sec"] = time_stamps[1]
    
    splited_audios = glob.glob(temp_outpt_wac + "*wav")
    
    # apply the text extraction in each splitted audio files 
    audio_sentences_timestamp = pd.DataFrame()
    for audio_clip_path in splited_audios:
        if filetype in audio_clip_path:
            extracted_txt = get_large_audio_transcription(path = audio_clip_path )
            temp = extracted_txt[1]
        else:
            pass
    #     audio_sentences_timestamp
        audio_sentences_timestamp = pd.concat([audio_sentences_timestamp, temp], axis = 0)
        
    # cleaned and wrnagel data intopandas df
    audio_sentences_timestamp["index"] = audio_sentences_timestamp.video_file.map(lambda x: int(x.split('_')[0]) )
    audio_sentences_timestamp = audio_sentences_timestamp.set_index("index")
    audio_sentences_timestamp = audio_sentences_timestamp.sort_index(ascending=True)
    
    
    # join with the time_stamps_df
    audio_sentences_timestamp = audio_sentences_timestamp.merge(time_stamps_df, left_index=True, right_index=True)
    audio_sentences_timestamp["audio_source"] = audio_sentences_timestamp.video_file.map(lambda x: re.sub(r'[0-9]', '', re.sub(r'[0-9]_', '', x)) )
    
    # clean divided audio clips from the folder
    for filePath in splited_audios:
        if os.path.exists(filePath):
            os.remove(filePath)
        else:
            print("Can not delete the file as it doesn't exists")
    
    return audio_sentences_timestamp                                                         



# audio_word_timestamp_df = pd.DataFrame()

# for audio_path in list_of_audios:
#     temp_awt = generate_timestamped_sentences(audio_path = audio_path, temp_outpt_wac = 'Audios/splitted_audio/', clip_length_sec= 10 )
#     audio_word_timestamp_df = pd.concat([audio_word_timestamp_df, temp_awt], axis = 0)
# audio_word_timestamp_df = audio_word_timestamp_df.reset_index(drop=True)

# print(audio_word_timestamp_df.shape)
# print(list(audio_word_timestamp_df['audio_source'].unique()) )

# # save this file into csv 
# audio_word_timestamp_df.to_csv("Targeted_Output/Phase1/0_raw_audio_word_timestamp.csv", index = False)