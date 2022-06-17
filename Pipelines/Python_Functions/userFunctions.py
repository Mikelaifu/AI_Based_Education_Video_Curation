# import en_core_web_sm
# nlp=en_core_web_sm.load()

def pptToText(p):
    from pptx import Presentation
    import pandas as pd
    #https://stackoverflow.com/questions/39418620/extracting-text-from-multiple-powerpoint-files-using-python
    txt = []
    prs = Presentation(p)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                txt.append(shape.text.strip())
    df = pd.DataFrame(txt, columns=['PPT Sub Text'])
    df['PPT File'] = p
    return df

def videoToAudio(video_file, output_ext='mp3'):
    import subprocess
    import os
    filename, ext = os.path.splitext(video_file)
    subprocess.call(['ffmpeg', '-y', '-i', video_file, f'{filename}.{output_ext}'], 
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.STDOUT)

def clientVideoJsonToDf(jsonFile):
    import pandas as pd
    import json
    with open(jsonFile,'r') as f:
        data = json.loads(f.read())
    df = pd.json_normalize(data, record_path=['alternatives'], meta=['resultEndTime'])
    df['Video ID'] = 'clientVideo'
    df.rename(columns={'transcript':'text'},inplace=True)
    return df

def youtubeTranscriptApiData(searchText,category,i):
    from pytube import Search
    import pandas as pd
    from youtube_transcript_api import YouTubeTranscriptApi
    transcript = []
    nTranscript = []

    #https://pytube.io/en/latest/user/search.html
    s = Search(searchText)
    vids = s.results

    [vids.append(Search(s.get_next_results).results) for k in range(1,i)]

    df = pd.DataFrame()

    for v in vids:
        try:
            #https://www.geeksforgeeks.org/python-downloading-captions-from-youtube/
            vidTranscipt = YouTubeTranscriptApi.get_transcript(v.video_id)
        except:
            nTranscript.append(v)
        else:
            transcript.append(v)
            sdf = pd.DataFrame(vidTranscipt)
            sdf['Video ID'] = v.video_id
            sdf['Title'] = v.title
            sdf['Age Restricted'] = v.age_restricted
            #sdf['Key Words'] = v.keywords
            sdf['Length (min)'] = (v.length/60)
            sdf['Views (thous)'] = (v.views/1000)
            sdf['Category'] = category
        df = pd.concat([df,sdf])
    print('Percentage of Videos w/Transcript: %s'%(len(transcript)*100/(len(transcript)+len(nTranscript)))+'%')
    return df

# https://realpython.com/natural-language-processing-spacy-python/
def is_token_allowed(token):
    if (not token or not token.text.strip() or token.is_stop or token.is_punct):
        return False
    return True

def preprocess_token(token):
    return token.lemma_.strip().lower()

def filteredTokens(txt):
    filteredTxt = ' '.join([preprocess_token(token) for token in txt if is_token_allowed(token)])
    return filteredTxt

def countTfidfWord2VecScore_between_2_text(new_doc_txt,reference_doc_txt):
    import en_core_web_lg
    nlp = en_core_web_lg.load()
    from sklearn.feature_extraction.text import CountVectorizer
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    import numpy as np
    
    reference_textClean = filteredTokens(nlp(reference_doc_txt))
    vid_text = filteredTokens(nlp(new_doc_txt))
    
    # measure cosine similarity from CountVectorizer
    vectorizer1 = CountVectorizer()
    X1 = vectorizer1.fit_transform([vid_text, reference_textClean])
    
    # measure cosine similarity from TfidfVectorizer
    vectorizer2 = TfidfVectorizer()
    X2 = vectorizer2.fit_transform([vid_text, reference_textClean])

    # measure cosine similarity from Word2Vec
    txt = [vid_text, reference_textClean]
    w2vEmbedding = [nlp(t.lower()).vector for t in txt]
    
    return [cosine_similarity(X1,X1)[0,1],cosine_similarity(X2,X2)[0,1],cosine_similarity(np.stack([w2vEmbedding[0],w2vEmbedding[1]]))[0,1]]
    
def countTfidfWord2VecScoreVideo(ppt,vid):
    import spacy
    import en_core_web_lg
    nlp = en_core_web_lg.load()
    import pandas as pd
    from sklearn.feature_extraction.text import CountVectorizer
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
    import numpy as np
    ppt = ppt.replace('Source: ebird.org', '').replace('Source: en.wikipedia.org, commons.wikimedia.org', '')
    pptClean = filteredTokens(nlp(ppt['PPT Sub Text'].to_string()))
    vidGroup = [vid[(vid['Video ID'] == x)]['text'] for x in vid['Video ID'].unique()]
    vidScore = []
    for v in vidGroup:
        v = v.replace('Source: ebird.org', '').replace('Source: en.wikipedia.org, commons.wikimedia.org', '')
        vidClean = filteredTokens(nlp(v.to_string()))
        vectorizer = CountVectorizer()
        X = vectorizer.fit_transform([pptClean,vidClean])
        vidScore.append(cosine_similarity(X,X)[0,1])
    countScoreDf = pd.concat([pd.DataFrame(vid['Video ID'].unique(), columns=['Video ID']), pd.DataFrame(vidScore, columns=['Count Vector'])], axis=1)
    vidScore = []
    for v in vidGroup:
        v = v.replace('Source: ebird.org', '').replace('Source: en.wikipedia.org, commons.wikimedia.org', '')
        vidClean = filteredTokens(nlp(v.to_string()))
        vectorizer = TfidfVectorizer()
        X = vectorizer.fit_transform([pptClean,vidClean])
        vidScore.append(cosine_similarity(X,X)[0,1])
    countTfidfScoreDf = pd.concat([countScoreDf, pd.DataFrame(vidScore, columns=['TFIDF Vector'])], axis=1)
    vidScore = []
    for v in vidGroup:
        v = v.replace('Source: ebird.org', '').replace('Source: en.wikipedia.org, commons.wikimedia.org', '')
        vidClean = filteredTokens(nlp(v.to_string()))
        txt = [pptClean,vidClean]
        w2vEmbedding = [nlp(t.lower()).vector for t in txt]
        vidScore.append(cosine_similarity(np.stack([w2vEmbedding[0],w2vEmbedding[1]]))[0,1])
    countTfidfWord2VecScoreDf = pd.concat([countTfidfScoreDf, pd.DataFrame(vidScore, columns=['Word To Vector'])], axis=1)
    countTfidfWord2VecScoreDf['Category'] = vid['Category'].iloc[0]
    countTfidfWord2VecScoreDf = countTfidfWord2VecScoreDf.sort_values(by= 'Word To Vector', ascending=False)
    return countTfidfWord2VecScoreDf